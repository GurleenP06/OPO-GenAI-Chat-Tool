import os
import pandas as pd
import numpy as np
import nltk
from pathlib import Path
from typing import List, Dict, Tuple, Optional
import logging
from dataclasses import dataclass
from tqdm import tqdm
import warnings

from pdfminer.high_level import extract_text as extract_pdf_text
from docx import Document
import pptx
import openpyxl

from sentence_transformers import SentenceTransformer
import faiss

import config

warnings.filterwarnings('ignore')

logging.basicConfig(
    level=logging.INFO,
    format='%(message)s'
)
logger = logging.getLogger(__name__)

logging.getLogger('pdfminer').setLevel(logging.WARNING)

nltk.data.path.append(r'C:/Users/e40082553/Desktop/CODES/OPO LLM Trial/models/nltk_data')

@dataclass
class ProcessedDocument:
    original_filename: str  # Store original filename (PDF, DOCX, etc.)
    text_filename: str      # Store text filename for reference
    text: str
    source_url: str
    chunks: List[str] = None


class DocumentProcessor:
    SUPPORTED_FORMATS = (".pdf", ".docx", ".pptx", ".xlsx", ".xls", ".xlsm")
    
    def __init__(self, 
                 input_dir: str,
                 output_dir: str,
                 chunk_size: int = 512,
                 chunk_overlap: int = 64,
                 batch_size: int = 32):
        
        self.input_dir = Path(input_dir)
        self.output_dir = Path(output_dir)
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.batch_size = batch_size
        
        self.text_output_dir = self.output_dir / "Text"
        self.text_output_dir.mkdir(parents=True, exist_ok=True)
        
        print(f"Loading embedding model: {config.EMBEDDING_MODEL_PATH}")
        self.embedding_model = SentenceTransformer(
            config.EMBEDDING_MODEL_PATH, 
            device=config.DEVICE
        )
        
        self.metadata_path = Path("./metadata.csv")
        self.faiss_index_path = Path("./faiss_index")
        self.source_links_path = self.input_dir / "SourceLinks.xlsx"
    
    def extract_text_from_pdf(self, pdf_path: Path) -> str:
        try:
            return extract_pdf_text(str(pdf_path))
        except Exception as e:
            logger.error(f"Error extracting text from PDF {pdf_path}: {e}")
            return ""
    
    def extract_text_from_docx(self, docx_path: Path) -> str:
        try:
            doc = Document(str(docx_path))
            return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
        except Exception as e:
            logger.error(f"Error extracting text from DOCX {docx_path}: {e}")
            return ""
    
    def extract_text_from_pptx(self, pptx_path: Path) -> str:
        try:
            presentation = pptx.Presentation(str(pptx_path))
            text = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        except Exception as e:
            logger.error(f"Error extracting text from PPTX {pptx_path}: {e}")
            return ""
    
    def extract_text_from_xlsx(self, xlsx_path: Path) -> str:
        try:
            wb = openpyxl.load_workbook(str(xlsx_path), data_only=True)
            text = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value:
                            text.append(str(cell.value))
            return "\n".join(text)
        except Exception as e:
            logger.error(f"Error extracting text from XLSX {xlsx_path}: {e}")
            return ""
    
    def extract_text_from_xls(self, xls_path: Path) -> str:
        try:
            sheets = pd.read_excel(str(xls_path), sheet_name=None, engine='xlrd')
            text_lines = []
            for sheet_name, df in sheets.items():
                text_lines.append(f"Sheet: {sheet_name}")
                for col in df.columns:
                    for val in df[col]:
                        if pd.notnull(val):
                            text_lines.append(str(val))
            return "\n".join(text_lines)
        except Exception as e:
            logger.error(f"Error extracting text from XLS {xls_path}: {e}")
            return ""
    
    def extract_text(self, file_path: Path) -> str:
        suffix = file_path.suffix.lower()
        
        extractors = {
            '.pdf': self.extract_text_from_pdf,
            '.docx': self.extract_text_from_docx,
            '.pptx': self.extract_text_from_pptx,
            '.xlsx': self.extract_text_from_xlsx,
            '.xlsm': self.extract_text_from_xlsx,
            '.xls': self.extract_text_from_xls
        }
        
        extractor = extractors.get(suffix)
        if extractor:
            return extractor(file_path)
        else:
            logger.warning(f"Unsupported file format: {suffix}")
            return ""
    
    def create_chunks(self, text: str) -> List[str]:
        words = nltk.word_tokenize(text)
        chunks = []
        
        if self.chunk_size > self.chunk_overlap:
            step = self.chunk_size - self.chunk_overlap
        else:
            step = self.chunk_size
        
        i = 0
        while i < len(words):
            chunk = words[i:i + self.chunk_size]
            chunks.append(' '.join(chunk))
            i += step
        
        return chunks
    
    def load_source_links(self) -> Dict[str, str]:
        if not self.source_links_path.exists():
            logger.warning(f"Source links file not found: {self.source_links_path}")
            return {}
        
        try:
            df = pd.read_excel(self.source_links_path)
            return dict(zip(df['File Name'], df['Original URL']))
        except Exception as e:
            logger.error(f"Error loading source links: {e}")
            return {}
    
    def process_documents(self) -> List[ProcessedDocument]:
        print("\nDOCUMENT PROCESSING")
        print("=" * 50)
        
        source_url_mapping = self.load_source_links()
        
        processed_docs = []
        
        files_to_process = []
        for root, _, files in os.walk(self.input_dir):
            if Path(root).name == "Text":
                continue
            
            for file in files:
                if file.lower().endswith(self.SUPPORTED_FORMATS):
                    files_to_process.append(Path(root) / file)
        
        total_files_in_dir = sum(1 for root, _, files in os.walk(self.input_dir) 
                                for file in files 
                                if not Path(root).name == "Text" and 
                                file.lower().endswith(self.SUPPORTED_FORMATS + (".xlsx",)))
        
        print(f"Found {len(files_to_process)} documents to process")
        print(f"(Total files in directory including SourceLinks: {total_files_in_dir})")
        
        print("\nExtracting text from documents...")
        success_count = 0
        failed_files = []
        
        for i, file_path in enumerate(files_to_process, 1):
            if i % 10 == 0:
                print(f"  Processed {i}/{len(files_to_process)} files...", end='\r')
            
            text = self.extract_text(file_path)
            if not text:
                failed_files.append(file_path.name)
                continue
            
            source_url = source_url_mapping.get(file_path.name, "")
            
            text_file_name = file_path.stem + ".txt"
            text_file_path = self.text_output_dir / text_file_name
            
            with open(text_file_path, "w", encoding="utf-8") as f:
                f.write(text)
            
            doc = ProcessedDocument(
                original_filename=file_path.name,  # Store original filename
                text_filename=text_file_name,      # Store text filename
                text=text,
                source_url=source_url
            )
            processed_docs.append(doc)
            success_count += 1
        
        print(f"\nSuccessfully processed {success_count} documents")
        if failed_files:
            print(f"Failed to extract text from {len(failed_files)} files:")
            for failed_file in failed_files[:5]:
                print(f"   - {failed_file}")
            if len(failed_files) > 5:
                print(f"   ... and {len(failed_files) - 5} more")
        
        return processed_docs
    
    def create_chunks_for_documents(self, documents: List[ProcessedDocument]) -> Tuple[List[str], List[Dict]]:
        print("\nCHUNKING")
        print("=" * 50)
        print(f"Creating chunks (size={self.chunk_size}, overlap={self.chunk_overlap})...")
        
        all_chunks = []
        all_metadata = []
        
        for doc in documents:
            chunks = self.create_chunks(doc.text)
            doc.chunks = chunks
            
            for chunk in chunks:
                all_chunks.append(chunk)
                all_metadata.append({
                    "chunk_text": chunk,
                    "original_filename": doc.original_filename,  # Store original filename
                    "text_filename": doc.text_filename,          # Keep text filename for reference
                    "source_url": doc.source_url
                })
        
        print(f"Created {len(all_chunks)} chunks from {len(documents)} documents")
        avg_chunks_per_doc = len(all_chunks) / len(documents) if documents else 0
        print(f"  Average chunks per document: {avg_chunks_per_doc:.1f}")
        
        return all_chunks, all_metadata
    
    def batch_encode(self, texts: List[str]) -> np.ndarray:
        embeddings = []
        total_batches = (len(texts) + self.batch_size - 1) // self.batch_size
        
        print(f"Generating embeddings for {len(texts)} chunks...")
        print(f"  Batch size: {self.batch_size}")
        print(f"  Total batches: {total_batches}")
        
        for i in range(0, len(texts), self.batch_size):
            batch_num = (i // self.batch_size) + 1
            if batch_num % 10 == 0 or batch_num == total_batches:
                print(f"  Processing batch {batch_num}/{total_batches}...", end='\r')
            
            batch = texts[i:i + self.batch_size]
            batch_embeddings = self.embedding_model.encode(
                batch, 
                convert_to_numpy=True,
                show_progress_bar=False
            )
            embeddings.append(batch_embeddings)
        
        print(f"\nGenerated embeddings for all chunks")
        return np.vstack(embeddings)
    
    def create_vector_database(self, chunks: List[str], metadata: List[Dict]):
        print("\nVECTOR DATABASE")
        print("=" * 50)
        
        embeddings = self.batch_encode(chunks)
        
        print("\nCreating FAISS index...")
        dimension = embeddings.shape[1]
        n_neighbors = min(32, len(chunks))  
        
        index = faiss.IndexHNSWFlat(dimension, n_neighbors)
        index.add(np.array(embeddings, dtype=np.float32))
        
        faiss.write_index(index, str(self.faiss_index_path))
        print(f"Saved FAISS index: {self.faiss_index_path}")
        
        metadata_df = pd.DataFrame(metadata)
        metadata_df.to_csv(self.metadata_path, index=False)
        print(f"Saved metadata: {self.metadata_path}")
    
    def run_pipeline(self):
        print("\n" + "="*60)
        print("DOCUMENT PROCESSING PIPELINE")
        print("="*60)
        print(f"\nConfiguration:")
        print(f"  Input directory: {self.input_dir}")
        print(f"  Output directory: {self.output_dir}")
        print(f"  Chunk size: {self.chunk_size} tokens")
        print(f"  Chunk overlap: {self.chunk_overlap} tokens")
        print(f"  Embedding model: {Path(config.EMBEDDING_MODEL_PATH).name}")
        print(f"  Device: {config.DEVICE}")
        
        documents = self.process_documents()
        
        if not documents:
            print("\nNo documents were processed successfully")
            return
        
        chunks, metadata = self.create_chunks_for_documents(documents)
        
        self.create_vector_database(chunks, metadata)
        
        print("\n" + "="*60)
        print("PIPELINE COMPLETED SUCCESSFULLY!")
        print("="*60)
        print(f"\nSummary:")
        print(f"  Documents processed: {len(documents)}")
        print(f"  Total chunks created: {len(chunks)}")
        print(f"  Output files:")
        print(f"     - Text files: {self.text_output_dir}")
        print(f"     - FAISS index: {self.faiss_index_path.absolute()}")
        print(f"     - Metadata: {self.metadata_path.absolute()}")
        print("\n" + "="*60)


def main():
    processor = DocumentProcessor(
        input_dir=config.TXT_DIRECTORY,
        output_dir=config.TXT_DIRECTORY,
        chunk_size=512,
        chunk_overlap=64,
        batch_size=32
    )
    
    processor.run_pipeline()


if __name__ == "__main__":
    main()
